"""
score.py
===========

Este script calcula uma pontuação quantitativa para os entregáveis do laboratório
com base na rubrica definida em `rubric.md`. A pontuação total é 100 pontos,
distribuídos entre várias categorias. O script verifica a existência de
artefatos, executa algumas validações básicas e atribui pontos parciais
conforme os critérios atendidos.

Uso:

```
python evaluation/score.py --root .
```

Por padrão, o diretório raiz do projeto é o diretório atual. Ajuste conforme
sua estrutura.
"""

import argparse
import os
import json
import csv

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None  # python-pptx pode não estar instalado


def evaluate_architecture(root: str) -> int:
    """Avalia a categoria Arquitetura & Design (30 pts)."""
    score = 0
    # Completude dos componentes (10 pts)
    required_dirs = ['infra', 'services', 'ml']
    found = sum(1 for d in required_dirs if os.path.isdir(os.path.join(root, d)))
    score += int((found / len(required_dirs)) * 10)
    # Redundância e failover (10 pts)
    redundancy_score = 0
    # Procura palavras-chave indicando redundância/failover em arquivos de infraestrutura
    keywords = ['autoscaler', 'replica', 'failover', 'multi_az', 'as_count']
    infra_path = os.path.join(root, 'infra')
    for dirpath, _, files in os.walk(infra_path):
        for fname in files:
            if fname.endswith(('.yaml', '.tf', '.hcl')):
                try:
                    with open(os.path.join(dirpath, fname), 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read().lower()
                        if any(k in content for k in keywords):
                            redundancy_score = 10
                            break
                except Exception:
                    continue
        if redundancy_score:
            break
    score += redundancy_score
    # Documentação SysML/UML (10 pts)
    blueprint = os.path.join(root, 'deliverable', '2_architecture_blueprint.pdf')
    if os.path.isfile(blueprint):
        score += 10
    return score


def evaluate_diagnosis(root: str) -> int:
    """Avalia a categoria Diagnóstico & Análise (20 pts)."""
    score = 0
    diag_path = os.path.join(root, 'deliverable', '1_diagnosis.pdf')
    if os.path.isfile(diag_path):
        # Presença do arquivo vale 10 pts
        score += 10
    # Segundo subcritério: uso de métricas e análise de rede. Verificamos se
    # o gerador de grafo e scripts de simulação existem.
    graph_gen = os.path.join(root, 'data', 'generate_graph.py')
    sim_script = os.path.join(root, 'sim', 'network_failure_sim.py')
    if os.path.isfile(graph_gen) and os.path.isfile(sim_script):
        score += 10
    return score


def evaluate_implementation(root: str) -> int:
    """Avalia a categoria Implementação & Observabilidade (20 pts)."""
    score = 0
    # Protótipo funcional (10 pts): verificamos se as pastas de serviços e mlflow existem
    # e se o CI passa (podemos checar se testes existem e são executáveis).
    services_ok = os.path.isdir(os.path.join(root, 'services'))
    tests_path = os.path.join(root, 'tests')
    if services_ok and os.path.isdir(tests_path):
        # Considera que CI é configurado se existe .github/workflows
        if os.path.isdir(os.path.join(root, '.github', 'workflows')):
            score += 10
    # Dashboards e alertas (10 pts): checa se há dashboards e alertas de Prometheus
    grafana_path = os.path.join(root, 'observability', 'grafana', 'dashboards')
    prom_alerts = os.path.join(root, 'observability', 'prometheus', 'alerts.yml')
    if os.path.isdir(grafana_path) and os.path.isfile(prom_alerts):
        score += 10
    return score


def evaluate_simulation(root: str) -> int:
    """Avalia a categoria Simulação & Quantitativo (15 pts)."""
    score = 0
    # Monte Carlo + interpretação (8 pts)
    report_csv = os.path.join(root, 'deliverable', '3_simulation_report.csv')
    if os.path.isfile(report_csv):
        try:
            with open(report_csv, newline='', encoding='utf-8', errors='ignore') as f:
                reader = csv.DictReader(f)
                # contar quantas combinações distintas de parâmetros
                p_nodes = set()
                for row in reader:
                    if 'p_node' in row:
                        p_nodes.add(row['p_node'])
                if len(p_nodes) >= 3:
                    score += 8
                else:
                    score += int((len(p_nodes) / 3) * 8)
        except Exception:
            pass
    # Propagação e stress (7 pts): verifica se scripts de simulação existem
    if os.path.isfile(os.path.join(root, 'sim', 'network_failure_sim.py')) and \
       os.path.isfile(os.path.join(root, 'sim', 'backpressure_sim.py')):
        score += 7
    return score


def evaluate_communication(root: str) -> int:
    """Avalia a categoria Comunicação & Estratégia (15 pts)."""
    score = 0
    # Deck executivo (8 pts): verifica existência de pptx e número de slides >=5
    deck_path = os.path.join(root, 'deliverable', '5_board_deck.pptx')
    if os.path.isfile(deck_path):
        pts = 0
        if Presentation is not None:
            try:
                prs = Presentation(deck_path)
                if len(prs.slides) >= 5:
                    pts = 8
                else:
                    pts = int((len(prs.slides) / 5) * 8)
            except Exception:
                pts = 4
        else:
            # Se não puder verificar slides, atribui metade dos pontos
            pts = 4
        score += pts
    # Runbooks e plano de mitigação (7 pts)
    responders = os.path.join(root, 'playbooks', 'runbook_responders.md')
    incidents = os.path.join(root, 'playbooks', 'incident_scenarios.md')
    postmortem = os.path.join(root, 'playbooks', 'postmortem_template.md')
    incident_template = os.path.join(root, 'playbooks', 'incident_template.md')
    present = sum(os.path.isfile(p) for p in [responders, incidents, postmortem, incident_template])
    score += int((present / 4) * 7)
    return score


def main():
    parser = argparse.ArgumentParser(description="Calcula a pontuação do laboratório Helius.")
    parser.add_argument('--root', default='.', help="Diretório raiz do projeto")
    args = parser.parse_args()
    root = os.path.abspath(args.root)

    scores = {}
    scores['architecture_design'] = evaluate_architecture(root)
    scores['diagnosis_analysis'] = evaluate_diagnosis(root)
    scores['implementation_observability'] = evaluate_implementation(root)
    scores['simulation_quantitative'] = evaluate_simulation(root)
    scores['communication_strategy'] = evaluate_communication(root)
    total = sum(scores.values())
    result = {
        'total_score': total,
        'breakdown': scores
    }
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()